#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Kontrollerer den byggede præsentation, før den sendes.

    python3 validate_deck.py Screening.pptx

`validate_data.py` tjekker tallene, før der bygges. Denne tjekker filen, der kom ud:
at pakken overhovedet kan åbnes af PowerPoint, og at slidesene overholder de regler,
skillen har sat sig — tre farver plus gråtoner, mindst 12 pt, intet uden for kanten,
og grafer som rigtige diagrammer frem for billeder.

Den kan ikke se, om en etiket er landet oven i en søjle. Kig altid på slidesene som
billeder bagefter.

Kræver `python-pptx` og `lxml`.

FEJL blokerer for levering. ADVARSEL er noget, du bør se på. Exitkode 1 ved mindst én fejl.
"""
import sys, zipfile, re, collections

try:
    from lxml import etree
    from pptx import Presentation
except ImportError:
    print("FEJL: python-pptx og lxml mangler. Kør:  pip install python-pptx lxml")
    sys.exit(2)

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
CH = "http://schemas.openxmlformats.org/drawingml/2006/chart"
PR = "http://schemas.openxmlformats.org/package/2006/relationships"

EMU = 914400
SLIDES = 6          # skillen leverer præcis seks slides
MIN_PT = 12.0       # mindste brødtekst
MAKS_FARVER = 3     # ud over gråtoner
TOL = EMU // 100    # 0,01 tomme at slå af på ved kanttjek

FEJL, ADVARSEL = [], []
def fejl(m): FEJL.append(m)
def advar(m): ADVARSEL.append(m)


def er_graa(v):
    """Gråtone = ingen nævneværdig mætning. Bedre end en fast liste over paletten,
    som skal vedligeholdes, hver gang der kommer en ny neutral tone til."""
    r, g, b = int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)
    return max(r, g, b) - min(r, g, b) <= 20


def tjek_pakken(sti):
    """Åbner zip'en og læser hver XML-del. En del, der ikke er velformet, giver en
    fil, PowerPoint nægter at åbne — og det opdager man først hos modtageren."""
    z = zipfile.ZipFile(sti)
    beskadiget = z.testzip()
    if beskadiget:
        fejl(f"beskadiget del i pakken: {beskadiget}")
    for navn in z.namelist():
        if navn.endswith((".xml", ".rels")):
            try:
                etree.fromstring(z.read(navn))
            except Exception as e:
                fejl(f"ugyldig XML i {navn}: {e}")

    # Hver relationship skal pege på en del, der faktisk ligger i pakken.
    dele = set(z.namelist())
    for navn in [n for n in z.namelist() if n.endswith(".rels")]:
        base = navn.rsplit("_rels/", 1)[0]
        for rel in etree.fromstring(z.read(navn)).iter(f"{{{PR}}}Relationship"):
            if rel.get("TargetMode") == "External":
                continue
            maal = rel.get("Target")
            sti_ = maal[1:] if maal.startswith("/") else (base + maal).replace("\\", "/")
            while "/../" in sti_:
                sti_ = re.sub(r"[^/]+/\.\./", "", sti_, count=1)
            if sti_.lstrip("/") not in dele:
                fejl(f"{navn}: relationship peger på {maal}, som ikke findes i pakken")


def tjek_diagram(xml, nr, farver):
    """Tjekker ét diagram: skriftstørrelser og at etiketterne står lovligt."""
    for e in xml.iter(f"{{{A}}}srgbClr"):
        v = (e.get("val") or "").upper()
        if v and not er_graa(v):
            farver[v] += 1

    # Skrift i et diagram sættes med defRPr, ikke rPr — akseetiketter, dataetiketter
    # og titel ligger alle i pPr/defRPr. Begge tjekkes, så ingen af dem slipper under.
    smaa = set()
    for tag in (f"{{{A}}}defRPr", f"{{{A}}}rPr"):
        for e in xml.iter(tag):
            sz = e.get("sz")
            if sz and int(sz) / 100.0 < MIN_PT:
                smaa.add(int(sz) / 100.0)
    for pt in sorted(smaa):
        fejl(f"slide {nr}: tekst i diagrammet er {pt:g} pt (mindst {MIN_PT:g} pt)")

    # På en stablet søjle må etiketten kun stå inde i søjlen. Står den uden for,
    # tegner PowerPoint den oven i næste segment. Meldes én gang pr. graf —
    # pptxgenjs skriver positionen både på grafen og på hver serie.
    for graf in xml.iter(f"{{{CH}}}barChart"):
        gr = graf.find(f"{{{CH}}}grouping")
        if gr is None or gr.get("val") != "stacked":
            continue
        ulovlige = {p.get("val") for p in graf.iter(f"{{{CH}}}dLblPos")
                    if p.get("val") not in ("ctr", "inEnd", "inBase")}
        for v in sorted(ulovlige):
            fejl(f"slide {nr}: dataLabelPosition '{v}' er ulovlig på stablede søjler "
                 f"— brug ctr, inEnd eller inBase")


def tjek_slides(sti, farver):
    prs = Presentation(sti)
    B, H = prs.slide_width, prs.slide_height
    print(f"{B/EMU:.3f} × {H/EMU:.3f} tommer · {len(prs.slides)} slides")

    if round(B / EMU, 2) != 13.33 or round(H / EMU, 2) != 7.5:
        fejl(f"slidestørrelsen er {B/EMU:.2f}×{H/EMU:.2f} tommer — sæt pres.layout til "
             f"LAYOUT_WIDE (13,333×7,5) før første slide")
    if len(prs.slides) != SLIDES:
        fejl(f"{len(prs.slides)} slides — skillen leverer præcis {SLIDES}")

    grafer = 0
    for nr, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if "PICTURE" in str(sh.shape_type):
                fejl(f"slide {nr}: '{sh.name}' er et billede — grafer skal laves med "
                     f"addChart, så modtageren kan klikke på en søjle og se tallet bag")
            if sh.has_chart:
                grafer += 1

            v, o, b, h = sh.left, sh.top, sh.width, sh.height
            if None not in (v, o, b, h):
                if v < -TOL or o < -TOL or v + b > B + TOL or o + h > H + TOL:
                    fejl(f"slide {nr}: '{sh.name}' rager ud over slidet "
                         f"({v/EMU:.2f},{o/EMU:.2f} {b/EMU:.2f}×{h/EMU:.2f} tommer)")

            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size is not None and r.font.size.pt < MIN_PT:
                            fejl(f"slide {nr}: «{r.text[:34]}» er {r.font.size.pt:g} pt. "
                                 f"Mindst {MIN_PT:g} pt — flyt teksten til regnearket i "
                                 f"stedet for at skrive den mindre")

        x = etree.fromstring(s.part.blob)
        for e in x.iter(f"{{{A}}}srgbClr"):
            v = (e.get("val") or "").upper()
            if v and not er_graa(v):
                farver[v] += 1

        for rel in s.part.rels.values():
            if "chart" in rel.reltype and not rel.is_external:
                tjek_diagram(etree.fromstring(rel.target_part.blob), nr, farver)

    if not grafer:
        advar("ingen diagrammer i præsentationen — er graferne endt som billeder?")
    return grafer


def main():
    if len(sys.argv) < 2:
        print("brug: validate_deck.py <deck.pptx>"); sys.exit(2)
    sti = sys.argv[1]

    farver = collections.Counter()
    tjek_pakken(sti)
    grafer = tjek_slides(sti, farver)

    if len(farver) > MAKS_FARVER:
        fejl(f"{len(farver)} kategorifarver ({', '.join(sorted(farver))}) — højst "
             f"{MAKS_FARVER} plus gråtoner. Flere farver er ikke længere en kode, "
             f"læseren kan holde i hovedet")

    print(f"{grafer} diagrammer · kategorifarver: "
          f"{', '.join(sorted(farver)) if farver else 'ingen'}")
    for m in ADVARSEL:
        print(f"ADVARSEL  {m}")
    for m in FEJL:
        print(f"FEJL      {m}")
    print(f"\n{len(FEJL)} fejl · {len(ADVARSEL)} advarsler")
    if not FEJL:
        print("Kig på slidesene som billeder, før du sender dem. Validatoren kan ikke se, "
              "om en etiket er landet oven i en søjle.")
    sys.exit(1 if FEJL else 0)


if __name__ == "__main__":
    main()
